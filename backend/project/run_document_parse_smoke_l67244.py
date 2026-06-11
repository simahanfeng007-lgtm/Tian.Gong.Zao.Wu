from __future__ import annotations

import tempfile
import importlib.util
from pathlib import Path

from tiangong_agent_runtime.document_parser import parse_document, should_route_to_document_parse
from tiangong_agent_runtime.adapters.readonly_file_adapter import read_file_adapter
from tiangong_agent_runtime.adapters.document_parse_adapter import document_parse_adapter
from tiangong_agent_runtime.tool_invocation import ToolInvocation
from tiangong_agent_runtime.turn_context import TurnContext
from tiangong_agent_runtime.plan_bridge import PlanBridge
from tiangong_agent_runtime.runtime_entry import RuntimeEntry
from tiangong_agent_runtime.plan_schema import validate_and_build_plan
from tiangong_agent_shell.tool_bridge import ToolBridge

DOCS_OPTIONAL_DEPS = {
    "python-docx": "docx",
    "openpyxl": "openpyxl",
    "python-pptx": "pptx",
    "pypdf": "pypdf",
}


def missing_docs_optional_deps() -> list[str]:
    missing: list[str] = []
    for package_name, module_name in DOCS_OPTIONAL_DEPS.items():
        if importlib.util.find_spec(module_name) is None:
            missing.append(package_name)
    return missing


def assert_true(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def no_raw_garble(text: str) -> bool:
    bad = ["PK\x03\x04", "%PDF-", "\ufffd\ufffd\ufffd", "\x00"]
    return not any(item in text for item in bad)


def make_docx(path: Path) -> None:
    from docx import Document  # type: ignore
    doc = Document()
    doc.add_heading("临渊者文档解析验收", 1)
    doc.add_paragraph("这是 DOCX 正文第一段，用于验证正文提取。")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "值"
    table.cell(1, 0).text = "版本"
    table.cell(1, 1).text = "L6.72.44"
    doc.save(path)


def make_xlsx(path: Path) -> None:
    import openpyxl  # type: ignore
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "验收"
    ws.append(["项目", "状态", "备注"])
    ws.append(["文档解析", "通过", "openpyxl 摘要"])
    wb.save(path)


def make_pptx(path: Path) -> None:
    from pptx import Presentation  # type: ignore
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "临渊者 PPTX 验收"
    slide.placeholders[1].text = "每页文本摘要应可提取。"
    prs.save(path)


def make_pdf(path: Path) -> None:
    from pypdf import PdfWriter  # type: ignore
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with path.open("wb") as f:
        writer.write(f)


def main() -> int:
    missing = missing_docs_optional_deps()
    if missing:
        print(
            "document_parse_smoke SKIP: optional docs extra missing: "
            + ", ".join(missing)
            + "; install backend/project/requirements-docs.txt to run DOCX/XLSX/PPTX/PDF fixtures."
        )
        return 0
    with tempfile.TemporaryDirectory(prefix="l67244_docparse_") as tmp:
        root = Path(tmp)
        txt = root / "编码测试.txt"
        txt.write_bytes("中文 UTF-8 文本\n第二行".encode("utf-8"))
        gb = root / "gb18030文本.txt"
        gb.write_bytes("中文 GB18030 文本".encode("gb18030"))
        csv_path = root / "表格.csv"
        csv_path.write_text("姓名,分数\n临渊者,100\n", encoding="utf-8")
        code = root / "demo.py"
        code.write_text("def answer():\n    return 'ok'\n", encoding="utf-8")
        docx = root / "验收.docx"
        xlsx = root / "验收.xlsx"
        pptx = root / "验收.pptx"
        pdf = root / "空白.pdf"
        binary = root / "图片.png"
        zip_file = root / "素材.zip"
        make_docx(docx)
        make_xlsx(xlsx)
        make_pptx(pptx)
        make_pdf(pdf)
        binary.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 40)
        zip_file.write_bytes(b"PK\x03\x04binary zip header only")

        for path in [txt, gb, csv_path, code, docx, xlsx, pptx, pdf, binary, zip_file]:
            assert_true(should_route_to_document_parse(path), f"should route: {path.suffix}")
            result = parse_document(path)
            summary = result.get("human_readable_summary", "")
            assert_true(result.get("raw_bytes_hidden") is True, f"raw bytes hidden: {path.name}")
            assert_true("【文档解析】" in summary, f"summary projection: {path.name}")
            assert_true(no_raw_garble(summary), f"no raw/garble in summary: {path.name}")

        assert_true("DOCX" in parse_document(docx)["human_readable_summary"], "docx summary")
        assert_true("工作表" in parse_document(xlsx)["human_readable_summary"], "xlsx sheets")
        assert_true("幻灯片" in parse_document(pptx)["human_readable_summary"], "pptx slides")
        assert_true("PDF" in parse_document(pdf)["human_readable_summary"], "pdf guard/summary")
        assert_true("image" in parse_document(binary)["human_readable_summary"].lower(), "image metadata guard")
        assert_true("archive" in parse_document(zip_file)["human_readable_summary"].lower(), "zip metadata guard")

        ctx = TurnContext.create("读 docx", workspace=root)
        inv = ToolInvocation("read_file", {"path": docx.name})
        rf = read_file_adapter(inv, ctx)
        assert_true(rf.ok, "read_file delegated OK")
        assert_true(rf.data.get("parser_tool") == "document_parse", "read_file delegates document_parse")
        assert_true("PK" not in rf.output_summary and "\ufffd" not in rf.output_summary, "read_file no raw bytes")

        dp = document_parse_adapter(ToolInvocation("document_parse", {"path": xlsx.name}), ctx)
        assert_true(dp.ok, "document_parse adapter OK")
        assert_true("工作表" in dp.output_summary, "document_parse adapter summary")

        hint = "[桌面端主机文件访问提示]\n- desktop_relative_path=."
        plan = PlanBridge().build_plan(f"帮我读一下桌面的 {docx.name}\n\n{hint}")
        assert_true(plan and plan[0].tool_name == "document_parse", "host file task routes docx to document_parse")

        runtime_result = RuntimeEntry().run_text(f"帮我读一下桌面的 {docx.name}\n\n{hint}", workspace=root)
        assert_true(runtime_result.projection.status == "ok", "RuntimeEntry document_parse run ok")
        assert_true(runtime_result.results and runtime_result.results[0].tool_name == "document_parse", "RuntimeEntry executes document_parse tool")
        assert_true("DOCX" in runtime_result.projection.summary and "PK" not in runtime_result.projection.summary, "RuntimeEntry summary is safe")

    print("L6.72.44 document_parse smoke passed")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
